from http.server import BaseHTTPRequestHandler
import json
import tempfile
import os
import requests
from io import BytesIO
import base64
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_bytes
import pypdf

class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            content_length = int(self.headers['Content-Length'])
            post_data = self.rfile.read(content_length)
            data = json.loads(post_data)
            
            # Validate input
            if 'pdf_url' not in data and 'pdf_data' not in data:
                self.send_error_response(400, "Missing 'pdf_url' or 'pdf_data' in request")
                return
            
            # Download PDF
            if 'pdf_url' in data:
                pdf_content = self.download_pdf(data['pdf_url'])
            else:
                # Handle base64 encoded PDF data
                pdf_content = base64.b64decode(data['pdf_data'])
            
            if not pdf_content:
                self.send_error_response(400, "Failed to get PDF content")
                return
            
            # Convert PDF to PPT
            ppt_buffer = self.convert_pdf_to_ppt(pdf_content)
            
            if not ppt_buffer:
                self.send_error_response(500, "Conversion failed")
                return
            
            # Send successful response
            self.send_response(200)
            self.send_header('Content-type', 'application/json')
            self.send_header('Access-Control-Allow-Origin', '*')
            self.end_headers()
            
            response = {
                "success": True,
                "message": "PDF converted to PPT successfully",
                "ppt_data": base64.b64encode(ppt_buffer.getvalue()).decode('utf-8'),
                "filename": "converted_presentation.pptx"
            }
            
            self.wfile.write(json.dumps(response).encode('utf-8'))
            
        except Exception as e:
            self.send_error_response(500, f"Internal server error: {str(e)}")
    
    def download_pdf(self, pdf_url):
        """Download PDF from URL"""
        try:
            response = requests.get(pdf_url, timeout=30)
            response.raise_for_status()
            
            # Verify it's a PDF
            content_type = response.headers.get('content-type', '').lower()
            if 'pdf' not in content_type:
                # Check magic number for PDF
                if not response.content.startswith(b'%PDF'):
                    return None
            
            return response.content
        except Exception as e:
            print(f"Download error: {e}")
            return None
    
    def convert_pdf_to_ppt(self, pdf_content):
        """Convert PDF content to PowerPoint presentation"""
        try:
            # Create a new presentation
            prs = Presentation()
            
            # Method 1: Convert PDF pages to images and add as slides
            images = self.pdf_to_images(pdf_content)
            
            if images:
                for i, image in enumerate(images):
                    # Create blank slide layout
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Add image to slide
                    img_stream = BytesIO()
                    image.save(img_stream, format='PNG')
                    img_stream.seek(0)
                    
                    # Add picture to slide (fill entire slide)
                    left = top = Inches(0)
                    slide.shapes.add_picture(img_stream, left, top, 
                                           width=prs.slide_width, 
                                           height=prs.slide_height)
            else:
                # Method 2: Extract text and create text-based slides
                self.create_text_slides_from_pdf(prs, pdf_content)
            
            # Save presentation to buffer
            ppt_buffer = BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            return ppt_buffer
            
        except Exception as e:
            print(f"Conversion error: {e}")
            return None
    
    def pdf_to_images(self, pdf_content, dpi=150):
        """Convert PDF pages to images"""
        try:
            images = convert_from_bytes(
                pdf_content, 
                dpi=dpi, 
                fmt='png',
                use_pdftocairo=True
            )
            return images
        except Exception as e:
            print(f"PDF to images error: {e}")
            return None
    
    def create_text_slides_from_pdf(self, prs, pdf_content):
        """Extract text from PDF and create slides"""
        try:
            pdf_file = BytesIO(pdf_content)
            pdf_reader = pypdf.PdfReader(pdf_file)
            
            for page_num in range(min(len(pdf_reader.pages), 10)):  # Limit to 10 pages
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                if text.strip():
                    # Use title slide layout
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title and content
                    title_shape = slide.shapes.title
                    content_shape = slide.placeholders[1]
                    
                    title_shape.text = f"Slide {page_num + 1}"
                    content_shape.text = text[:1000]  # Limit text length
                    
        except Exception as e:
            print(f"Text extraction error: {e}")
    
    def send_error_response(self, code, message):
        """Send error response"""
        self.send_response(code)
        self.send_header('Content-type', 'application/json')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        
        error_response = {
            "success": False,
            "error": message,
            "code": code
        }
        
        self.wfile.write(json.dumps(error_response).encode('utf-8'))
    
    def do_OPTIONS(self):
        """Handle preflight requests"""
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()