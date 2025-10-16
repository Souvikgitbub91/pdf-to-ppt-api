from flask import Flask, request, jsonify
import requests
from io import BytesIO
import base64
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_bytes
import os

app = Flask(__name__)

@app.route('/')
def home():
    return jsonify({
        "message": "PDF to PPT API is running",
        "status": "active",
        "endpoints": {
            "GET /status": "Health check",
            "POST /convert": "Convert PDF to PPT"
        }
    })

@app.route('/status', methods=['GET'])
def status():
    return jsonify({
        "message": "PDF to PPT API is running",
        "status": "active",
        "version": "1.0"
    })

@app.route('/convert', methods=['POST'])
def convert_pdf_to_ppt():
    try:
        data = request.get_json()
        
        if not data or 'pdf_url' not in data:
            return jsonify({
                "success": False,
                "error": "Missing 'pdf_url' in request body"
            }), 400
        
        pdf_url = data['pdf_url']
        
        # Download PDF
        pdf_content = download_pdf(pdf_url)
        if not pdf_content:
            return jsonify({
                "success": False,
                "error": "Failed to download PDF from URL"
            }), 400
        
        # Convert to images and create PPT
        ppt_buffer = create_ppt_from_pdf(pdf_content)
        
        if not ppt_buffer:
            return jsonify({
                "success": False,
                "error": "Failed to create PowerPoint presentation"
            }), 500
        
        # Return success with base64 encoded PPT
        return jsonify({
            "success": True,
            "message": "PDF converted to PPT successfully",
            "ppt_data": base64.b64encode(ppt_buffer.getvalue()).decode('utf-8'),
            "filename": "converted_presentation.pptx",
            "pages": len(convert_from_bytes(pdf_content))
        })
        
    except Exception as e:
        return jsonify({
            "success": False,
            "error": f"Conversion error: {str(e)}"
        }), 500

def download_pdf(pdf_url):
    """Download PDF from URL"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        response = requests.get(pdf_url, headers=headers, timeout=30)
        response.raise_for_status()
        
        # Verify it's a PDF
        if response.content.startswith(b'%PDF'):
            return response.content
        return None
        
    except Exception as e:
        print(f"Download error: {e}")
        return None

def create_ppt_from_pdf(pdf_content):
    """Convert PDF to PowerPoint presentation"""
    try:
        # Create presentation
        prs = Presentation()
        
        # Convert PDF pages to images
        images = convert_from_bytes(pdf_content, dpi=150, fmt='PNG')
        
        for i, image in enumerate(images):
            # Create blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Convert image to bytes
            img_bytes = BytesIO()
            image.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            
            # Add image to slide (full screen)
            left = top = Inches(0)
            pic = slide.shapes.add_picture(img_bytes, left, top, 
                                         width=prs.slide_width, 
                                         height=prs.slide_height)
        
        # Save to buffer
        ppt_buffer = BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        return ppt_buffer
        
    except Exception as e:
        print(f"PPT creation error: {e}")
        return None

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)